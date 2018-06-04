import six
from pptx import Presentation
from pptx.util import Inches, Pt, Mm
import pandas as pd
from math import *

round_to_n = lambda x, n: round(x, -int(floor(log10(abs(x)))) + (n - 1))

"""
Source: https://github.com/robintw/PandasToPowerpoint
"""


def _do_formatting(value, format_str):
    """Format value according to format_str, and deal
    sensibly with format_str if it is missing or invalid."""
    if format_str == '':
        if type(value) in six.integer_types:
            format_str = ','
        elif type(value) is float:
            format_str = 'f'
        elif type(value) is str:
            format_str = 's'
    elif format_str[0] == '.':
        if format_str.endswith('R'):
            if type(value) in six.integer_types:
                value = round_to_n(value, int(format_str[1]))
                format_str = ','
        if not format_str.endswith('G'):
            format_str = format_str + "G"
    try:
        value = format(value, format_str)
    except:
        value = format(value, '')

    return value


def df_to_table(slide, df, left, top, width, height, colnames=None, col_formatters=None, include_index=True, units='mm'):
    """Converts a Pandas DataFrame to a PowerPoint table on the given
    Slide of a PowerPoint presentation.

    The table is a standard Powerpoint table, and can easily be modified with the Powerpoint tools,
    for example: resizing columns, changing formatting etc.

    Arguments:
     - slide: slide object from the python-pptx library containing the slide on which you want the table to appear
     - df: Pandas DataFrame with the data

    Optional arguments:
     - col_formatters: A n_columns element long list containing format specifications for each column.
     For example ['', ',', '.2'] does no special formatting for the first column, uses commas as thousands separators
     in the second column, and formats the third column as a float with 2 decimal places.
     - units: Default units are mm, but other options are inches, pt
     """
    if units.lower() == 'mm':
        unit_cls = Mm
    elif units.lower() == 'inches':
        unit_cls = Inches
    elif units.lower() == 'pt':
        unit_cls = Pt
    else:
        raise NotImplementedError("Unrecognised units %s" % units)

    rows, cols = df.shape
    nrows = rows + 1
    if include_index:
        ncols = cols + 1
    else:
        ncols = cols

    res = slide.shapes.add_table(
        nrows,
        ncols,
        unit_cls(left),
        unit_cls(top),
        unit_cls(width),
        unit_cls(height)
    )

    if colnames is None:
        colnames = list(df.columns)

    # Insert the column names
    # Cast to string just in case this is required
    for col_index, col_name in enumerate(colnames):
        if include_index:
            col_index += 1
        res.table.cell(0, col_index).text = str(col_name)

    # insert the row names if requested
    if include_index:
        for row_index, row_name in enumerate(list(df.index)):
            res.table.cell(row_index + 1, 0).text = row_name

    m = df.as_matrix()

    for row in range(rows):
        for col in range(cols):
            colidx = col + 1 if include_index else col
            val = m[row, col]

            if col_formatters is None:
                text = str(val)
            else:
                text = _do_formatting(val, col_formatters[col])

            res.table.cell(row + 1, colidx).text = text
            # res.table.cell(row+1, col).text_frame.fit_text()


def df_to_powerpoint(filename, df, **kwargs):
    """Converts a Pandas DataFrame to a table in a new, blank PowerPoint presentation.

    Creates a new PowerPoint presentation with the given filename, with a single slide
    containing a single table with the Pandas DataFrame data in it.

    The table is a standard Powerpoint table, and can easily be modified with the Powerpoint tools,
    for example: resizing columns, changing formatting etc.

    Arguments:
     - filename: Filename to save the PowerPoint presentation as
     - df: Pandas DataFrame with the data

     Kwargs:
     NB. All units are in millimetres, unless a kwarg `unit` is supplied
     - left: leftmost coordinate
     - top: margin from top
     - width, height: hopefully self-explanatory!

    All other arguments that can be taken by df_to_table() (such as col_formatters or rounding) can also
    be passed here.
    """
    kwargs.setdefault('left', 25)
    kwargs.setdefault('top', 25)
    kwargs.setdefault('width', 200)
    kwargs.setdefault('height', 150)

    if df.index.dtype != 'O':
        df = df.copy()
        df.index = df.index.astype(str)

    pres = Presentation()
    blank_slide_layout = pres.slide_layouts[6]
    slide = pres.slides.add_slide(blank_slide_layout)
    df_to_table(slide, df, **kwargs)
    pres.save(filename)